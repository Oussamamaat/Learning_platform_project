"""
Document Ingestion Service
─────────────────────────
Reads documents → strips formatting → chunks → embeds → stores in pgvector.

Supports .txt, .md, .pdf, .docx, .pptx, .xlsx, .csv, and images (via OCR --
see app.services.ocr). Every format is normalized to markdown BEFORE
chunking -- see parse_document_to_markdown -- so a single heading-aware
chunker (chunk_document) and a single citation-preservation guarantee
cover all of them, rather than one parser per format needing its own
chunking logic.

Usage:
    python -m app.services.ingestion --source_dir ./raw/shared --tenant_id company_abc
"""

import csv
import html
import io
import logging
import re
import uuid
import json
from functools import lru_cache
from pathlib import Path
from typing import Callable, Optional

import psycopg2
import psycopg2.extensions
import psycopg2.pool
import threading
from psycopg2.extras import execute_values
from sentence_transformers import SentenceTransformer
from langchain_text_splitters import RecursiveCharacterTextSplitter

from app.config import get_settings

logger = logging.getLogger(__name__)


# Reads settings.embedding_model at import time (not hardcoded) so
# scripts/migrate_to_bge_m3.py's rollback path (set EMBEDDING_MODEL back to
# the MiniLM name via env/.env) works without a code change. Does NOT
# affect app.services.generate_training_data.DEDUP_MODEL, which is a
# separate, deliberately still-MiniLM dataset-dedup threshold -- see
# app/config.py's embedding_model docstring.
DEFAULT_EMBEDDING_MODEL = get_settings().embedding_model
# 2000/250 (bge-m3, 8192-token window) as of 2026-08-13 -- was 400/50,
# sized against the OLD embedding model's actual 128-token max_seq_length
# (paraphrase-multilingual-MiniLM-L12-v2's own sentence_bert_config.json),
# not an arbitrary default. Raising this without also raising
# max_context_length (app/services/retrieval.py) and setting Ollama's
# num_ctx (app/services/llm.py) silently truncates citations -- see
# retrieval._build_context's docstring-adjacent comment for the concrete
# failure mode this caused.
CHUNK_SIZE = 2000
CHUNK_OVERLAP = 250
# Kept as a module constant for the callers/tests that import it; the
# EMBEDDING batch size now comes from settings.embedding_batch_size (see
# embed_chunks, which halves and retries on CUDA OOM) because that one has
# to be tunable per deployment. This value is the psycopg2 execute_values
# page size -- how many rows go in one INSERT statement, a pure
# network/parse-overhead knob with no GPU involvement.
BATCH_SIZE = 64
INSERT_PAGE_SIZE = 200

# raw/shared/<domain>/text/*.md -> "securite_physique" is the folder name,
# but "securite" is the Domain enum value used everywhere else in the app
# (app/models/schemas.py). Pre-existing mismatch, noted where it first bit
# (app/services/domain_context.py) -- normalized once, here, at ingest time.
DOMAIN_DIR_ALIASES = {"securite_physique": "securite"}

_HEADING_RE = re.compile(r"^(#{1,6})\s+(.*)$", re.MULTILINE)

# Legacy binary Office formats: rejected outright rather than attempted --
# no parser here reads the old OLE-based .doc/.ppt/.xls container format,
# and silently misreading one as UTF-8 text would be worse than a clear
# "convert first" error.
LEGACY_BINARY_FORMATS = {".doc": ".docx", ".ppt": ".pptx", ".xls": ".xlsx"}
_IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".tiff", ".tif"}
SUPPORTED_EXTENSIONS = {".md", ".txt", ".pdf", ".docx", ".pptx", ".xlsx", ".csv"} | _IMAGE_EXTENSIONS

# Rasterization DPI for pages that app.services.pdf_classify.classify_page
# sends to OCR (app.services.ocr). ~200 DPI balances legibility for small
# print (this corpus's article numbers/footnotes) against OCR call latency.
OCR_RENDER_DPI = 200
# Hard cap on rows rendered from one spreadsheet sheet -- silent truncation
# is unacceptable in a citation system, so a sheet over this size gets an
# explicit "[... N more rows omitted ...]" marker rather than either
# hanging on a huge sheet or dropping rows with no trace.
XLSX_MAX_ROWS = 5000


def _parse_text(file_path: Path) -> str:
    return file_path.read_text(encoding="utf-8")


def _parse_docx(file_path: Path) -> str:
    """`.docx` -> markdown. Word's own heading styles ("Heading 1".."Heading
    6", "Title") map directly to `#`.."######" so split_by_headings sees
    real document structure instead of a flat, unheaded text dump -- this
    is what lets the citation-preservation guarantee (chunk_document's
    docstring) extend to uploaded Word documents for free.
    """
    import docx

    document = docx.Document(str(file_path))
    lines: list[str] = []
    for para in document.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        style_name = (para.style.name or "").strip() if para.style else ""
        heading_match = re.match(r"^Heading (\d)$", style_name)
        if style_name == "Title":
            lines.append(f"# {text}")
        elif heading_match:
            level = min(int(heading_match.group(1)), 6)
            lines.append(f"{'#' * level} {text}")
        else:
            lines.append(text)

    for table in document.tables:
        lines.append("")
        for row in table.rows:
            cells = [c.text.strip().replace("\n", " ") for c in row.cells]
            lines.append(_render_row(cells))

    return "\n\n".join(lines)


def _parse_pptx(file_path: Path) -> str:
    """`.pptx` -> markdown. Slide title -> `## `, other text-frame shapes
    as body paragraphs, table shapes as markdown rows, speaker notes
    appended as a blockquote (kept for embedding -- notes often carry the
    actual explanatory content a slide's bullet points only gesture at).

    Table shapes are a SEPARATE shape type from the text-frame shapes the
    loop below already handled (`shape.has_table`, not `has_text_frame` --
    python-pptx never sets both on the same shape), so a slide's table was
    previously invisible to this parser entirely: silently dropped, no
    error, no record, indistinguishable from a slide that never had one.
    """
    import pptx

    presentation = pptx.Presentation(str(file_path))
    sections: list[str] = []
    for i, slide in enumerate(presentation.slides, start=1):
        title = None
        body_lines: list[str] = []
        for shape in slide.shapes:
            if shape.has_table:
                body_lines.append("")
                for row in shape.table.rows:
                    cells = [c.text.strip().replace("\n", " ") for c in row.cells]
                    body_lines.append(_render_row(cells))
                continue
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            if shape == slide.shapes.title:
                title = text
            else:
                body_lines.append(text)

        heading = f"## {title}" if title else f"## Slide {i}"
        section = [heading]
        section.extend(body_lines)

        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                quoted = "\n".join(f"> {line}" for line in notes.splitlines())
                section.append(quoted)

        sections.append("\n\n".join(section))

    return "\n\n".join(sections)


def _render_row(cells: list[str]) -> str:
    return "| " + " | ".join(c.replace("|", "\\|") for c in cells) + " |"


def _parse_xlsx(file_path: Path) -> str:
    """`.xlsx` -> markdown. `read_only=True, data_only=True` streams the
    file rather than loading the whole workbook, and reads each formula
    cell's last-cached value instead of the formula text. One `## <sheet
    name>` section per sheet, rows rendered as a markdown table, capped at
    XLSX_MAX_ROWS with an explicit truncation note.
    """
    import openpyxl

    workbook = openpyxl.load_workbook(str(file_path), read_only=True, data_only=True)
    sections: list[str] = []
    try:
        for sheet in workbook.worksheets:
            rows_rendered = 0
            lines = [f"## {sheet.title}"]
            for row in sheet.iter_rows(values_only=True):
                if all(cell is None for cell in row):
                    continue
                if rows_rendered >= XLSX_MAX_ROWS:
                    lines.append(f"[... remaining rows omitted after {XLSX_MAX_ROWS} ...]")
                    break
                cells = ["" if c is None else str(c) for c in row]
                lines.append(_render_row(cells))
                rows_rendered += 1
            if rows_rendered:
                sections.append("\n".join(lines))
    finally:
        workbook.close()

    return "\n\n".join(sections)


def _parse_csv(file_path: Path) -> str:
    """`.csv` -> markdown. `utf-8-sig` strips a BOM if present; delimiter
    sniffed rather than assumed comma, since this corpus's genre (exported
    regulatory/compliance spreadsheets) commonly uses `;` in French locale
    exports.
    """
    raw = file_path.read_bytes().decode("utf-8-sig")
    try:
        dialect = csv.Sniffer().sniff(raw[:2048])
    except csv.Error:
        dialect = csv.excel

    reader = csv.reader(io.StringIO(raw), dialect)
    lines = [f"# {file_path.stem}"]
    for i, row in enumerate(reader):
        if not any(cell.strip() for cell in row):
            continue
        if i >= XLSX_MAX_ROWS:
            lines.append(f"[... remaining rows omitted after {XLSX_MAX_ROWS} ...]")
            break
        lines.append(_render_row(row))

    return "\n".join(lines)


def _parse_image(file_path: Path) -> str:
    from app.services.ocr import get_ocr_engine

    engine = get_ocr_engine()
    return engine.image_to_markdown(file_path.read_bytes())


def _parse_pdf(
    file_path: Path,
    unprocessed_pages: Optional[list] = None,
    on_page_processed: Optional[Callable[[int, int], None]] = None,
    page_stats: Optional[dict] = None,
) -> str:
    """`.pdf` -> markdown. The text-vs-OCR decision is made PER PAGE, not
    per document: a digitally authored regulation with a scanned annex
    appended is the norm in this corpus's genre, and a document-level "does
    it have a text layer?" check would silently drop every scanned page of
    an otherwise-digital file. Each page is classified by
    app.services.pdf_classify.classify_page (font/image/density signals
    measured against real Moroccan PDFs, not a raw character-count floor --
    see that module's docstring for why) into one of four strategies:

      NATIVE        -- use the embedded text as-is.
      EMPTY         -- nothing readable; emit an empty section, never call
                       OCR (a blank page must not fail an otherwise fully
                       readable document).
      OCR_PREFERRED -- embedded text exists but is degraded (e.g. a diagram
                       whose layout the raw text order loses, or a table
                       carried by an embedded image the text layer never
                       had -- see pdf_classify's LARGE_IMAGE_MIN_MEGAPIXELS).
                       Native text and OCR output are MERGED, not one
                       replacing the other: measured on a real page, a
                       numeric constant existed ONLY in the native layer
                       while a table existed ONLY in the OCR output, so
                       trusting either alone loses real content. If OCR
                       returns nothing (empty/whitespace) or raises
                       OcrUnavailableError, the native text is kept as-is
                       rather than destroyed.
      OCR_REQUIRED  -- embedded text is unusable; if OCR is unavailable,
                       this ONE page is skipped (recorded into
                       `unprocessed_pages`, if given, as {"page", "reason"})
                       and parsing continues -- a real 80-page administrative
                       guide measured 4 OCR_REQUIRED pages against 66 NATIVE
                       ones, and failing the whole document over 5% of its
                       pages would discard the other 95% for no reason. The
                       one exception: if EVERY page that could contribute
                       text failed this way and nothing else in the document
                       produced any text either, that is not a partial
                       document, it is an unreadable one -- the original
                       OcrUnavailableError is re-raised so
                       app.services.ingest_jobs.process_source_file can
                       still fail loudly and actionably in that case.

    Every page that ends up with NO text at all -- classified EMPTY, or a
    successful-but-empty OCR call on an OCR_REQUIRED/OCR_PREFERRED page --
    is recorded into `unprocessed_pages`. Silently dropping a page's
    content while the document overall reports success (status='ready')
    is exactly the failure this exists to prevent: a 12-page loss inside
    an 80-page document was once invisible because only the OCR_REQUIRED
    failure path recorded anything.

    Emits `## Page N` headings so page numbers survive into chunks as a
    citation affordance.

    `page_stats`, if given, is filled in with {"page_count", "ocr_pages"}
    -- the counts this pass ALREADY computes while classifying each page.
    app.services.ingest_jobs._infer_parser_metadata needs exactly those two
    numbers to record the pdf_text/pdf_mixed/pdf_ocr audit trail, and used
    to get them by re-opening the file and running extract_text() +
    classify_page over every page a SECOND time, after ingestion had
    finished. Its docstring called that "effectively zero cost", which is
    true of classify_page but not of the extract_text() it feeds: text
    extraction is the dominant per-page cost on a native PDF, so an
    80-page document was paying for two full extraction passes. Handing
    the numbers forward removes the second one, and removes the chance of
    the two passes disagreeing.
    """
    from pypdf import PdfReader

    from app.services.ocr import OcrUnavailableError
    from app.services.pdf_classify import PageStrategy, classify_page

    reader = PdfReader(str(file_path))
    total_pages = len(reader.pages)
    sections: list[str] = []
    local_unprocessed: list[dict] = []
    any_text = False
    last_error: Optional[OcrUnavailableError] = None
    ocr_pages = 0

    for page_num, page in enumerate(reader.pages, start=1):
        text = (page.extract_text() or "").strip()
        decision = classify_page(page, text=text)
        page_recorded = False  # local_unprocessed already has an entry for this page

        if decision.strategy == PageStrategy.EMPTY:
            text = ""
        elif decision.strategy == PageStrategy.OCR_REQUIRED:
            try:
                # "heavy": this page has no usable native text at all, so
                # there is nothing to fall back on if OCR under-reads it.
                text = _ocr_pdf_page(file_path, page_num - 1, tier="heavy").strip()
            except OcrUnavailableError as e:
                # `e` -- the OCR call's OWN failure reason (subprocess error,
                # timeout, missing engine) -- not decision.reason (why this
                # page was classified OCR_REQUIRED in the first place). Those
                # are two different questions; logging decision.reason here
                # made every OCR failure look identical to a classification
                # note ("no font resources"), with no way to tell "OCR is
                # off" from "OCR ran and crashed" from a log line alone.
                logger.warning(
                    "page=%d of %s: OCR required (%s), but the OCR attempt "
                    "itself failed: %s -- page skipped, rest of document "
                    "continues",
                    page_num, file_path.name, decision.reason, e,
                )
                local_unprocessed.append({
                    "page": page_num, "reason": "ocr_required", "detail": str(e),
                })
                page_recorded = True
                last_error = e
                text = ""
        elif decision.strategy == PageStrategy.OCR_PREFERRED:
            native_text = text
            try:
                # "light": native_text below is merged in regardless, so
                # OCR here is only being asked for the embedded table or
                # figure -- the case the light engine measurably handles
                # (p15's CAD table, 4/5). It self-escalates to the heavy
                # engine when its output looks numerically empty.
                ocr_text = _ocr_pdf_page(file_path, page_num - 1, tier="light").strip()
                # Merge, don't replace: verified necessary against a real
                # page where a numeric constant existed ONLY in native_text
                # and a table existed ONLY in ocr_text. An empty ocr_text
                # (OCR ran, found nothing new) keeps native_text intact
                # rather than overwriting a good extraction with nothing.
                text = f"{native_text}\n\n{ocr_text}" if (native_text and ocr_text) else (ocr_text or native_text)
            except OcrUnavailableError:
                logger.info(
                    "page=%d of %s: OCR preferred but unavailable (%s) -- "
                    "falling back to embedded text",
                    page_num, file_path.name, decision.reason,
                )
                text = native_text
        # else NATIVE: text is already the embedded extraction.

        if text:
            any_text = True
        elif not page_recorded:
            # Every page that ends up with literally nothing -- classified
            # EMPTY, or a successful-but-empty OCR call above -- must be
            # recorded, not just the OCR_REQUIRED-failure case above. This
            # is what promotes source_files.status to 'partial' instead of
            # a silent 'ready' with content quietly missing.
            reason = "empty_page" if decision.strategy == PageStrategy.EMPTY else "no_text_extracted"
            local_unprocessed.append({"page": page_num, "reason": reason})

        if decision.strategy != PageStrategy.NATIVE:
            ocr_pages += 1
            logger.info(
                "page=%d of %s: %s -> %s",
                page_num, file_path.name, decision.strategy.value, decision.reason,
            )

        sections.append(f"## Page {page_num}\n\n{text}")

        if on_page_processed is not None:
            try:
                on_page_processed(page_num, total_pages)
            except Exception:
                # A progress callback is a side channel, not part of the
                # parse contract -- a bug in it (e.g. a DB hiccup writing
                # pages_done) must never fail an otherwise-successful
                # parse. Logged, not silently swallowed.
                logger.exception(
                    "on_page_processed callback failed for page=%d of %s (ignored)",
                    page_num, file_path.name,
                )

    if unprocessed_pages is not None:
        unprocessed_pages.extend(local_unprocessed)

    if page_stats is not None:
        page_stats["page_count"] = total_pages
        page_stats["ocr_pages"] = ocr_pages

    if not any_text and last_error is not None:
        # Every page that could have produced text failed identically --
        # this isn't a partial document, it's an entirely unreadable one.
        # Fail loudly rather than silently ingest something empty that
        # would look like a successful upload with nothing retrievable.
        # `unprocessed_pages` is populated above, before this raise, so a
        # caller inspecting it after catching this exception still sees
        # every page that was attempted, not an empty list.
        raise last_error

    return "\n\n".join(sections)


# A decimal (12.5) or an integer of 4+ digits (1970, 86400, 500000).
# Calibrated against measured output, see _classic_ocr_looks_incomplete.
_NUMERIC_EVIDENCE_RE = re.compile(r"\d+\.\d+|\d{4,}")


def _classic_ocr_looks_incomplete(text: str) -> bool:
    """Should this page's light-engine result be re-run on the heavy engine?

    The light engine (PP-OCRv5) reads running text and table cells well but
    measurably cannot read this corpus's formulas: on arabic_test.pdf p51 it
    scored 0/4 on the DPF constants that PaddleOCR-VL recovers 4/4. Escalating
    only the pages it likely failed keeps ~10x the speed on the rest.

    The signal, chosen because it separates the two pages where the light
    engine's real behaviour is known rather than because it sounds
    plausible:

        page                        decimals   >=4-digit ints   verdict
        p51 (formulas, scored 0/4)      0             0         escalate
        p15 (CAD table, scored 4/5)     0             1 (1970)  keep

    So: no decimal AND no long integer => escalate. This corpus's tables and
    formulas are numeric throughout (coordinates, thresholds, constants,
    years, zone areas), so a page whose OCR yields neither is one the light
    engine probably could not read.

    Two deliberate properties:
      - It errs TOWARD the heavy engine. A false escalation costs ~47s of
        GPU time; a false negative silently loses a formula from a
        regulatory document, which is the failure this whole change exists
        to prevent.
      - Empty/whitespace output escalates too, via the same rule.

    Calibrated on n=2 pages. If a real tenant document is later found where
    the light engine drops content this check misses, set
    settings.ocr_two_tier=False rather than widening the heuristic blindly.
    """
    return not _NUMERIC_EVIDENCE_RE.search(text or "")


def _ocr_pdf_page(file_path: Path, page_index: int, *, tier: str = "heavy") -> str:
    """Render one page and OCR it. `tier` selects the engine:

      "heavy" -- settings.ocr_paddle_engine. Used for OCR_REQUIRED pages,
                 which have NO usable native text to fall back on, so a
                 miss is unrecoverable.
      "light" -- settings.ocr_light_engine, ~10x faster (measured 5.5s vs
                 52.5s/page warm). Used for OCR_PREFERRED pages, whose
                 native text _parse_pdf is already merging in, so OCR is
                 only being asked for the embedded table/figure. Escalates
                 to "heavy" when _classic_ocr_looks_incomplete flags the
                 result.

    Rendering (open + render + PNG encode) measures ~200ms against a
    52,500ms OCR call, i.e. 0.4% -- so this deliberately keeps the simple
    open-per-page form rather than caching the PdfDocument across pages.
    Opening the document alone measured 1ms.
    """
    import pypdfium2 as pdfium

    from app.services.ocr import get_ocr_engine

    pdf = pdfium.PdfDocument(str(file_path))
    try:
        page = pdf[page_index]
        scale = OCR_RENDER_DPI / 72
        bitmap = page.render(scale=scale)
        pil_image = bitmap.to_pil()
        buf = io.BytesIO()
        pil_image.save(buf, format="PNG")
        image_bytes = buf.getvalue()
    finally:
        pdf.close()

    engine = get_ocr_engine()
    if tier != "light" or not get_settings().ocr_two_tier:
        return engine.image_to_markdown(image_bytes, tier="heavy")

    text = engine.image_to_markdown(image_bytes, tier="light")
    if _classic_ocr_looks_incomplete(text):
        # 1-based page number, matching _parse_pdf's own "page=%d" lines --
        # this log sits directly between two of them, and reporting the
        # 0-based index here made the same page look like two different
        # pages when reading the log top to bottom.
        logger.info(
            "page=%d of %s: light OCR returned no numeric evidence "
            "(%d chars) -- escalating to the heavy engine",
            page_index + 1, file_path.name, len(text or ""),
        )
        return engine.image_to_markdown(image_bytes, tier="heavy")
    return text


_PARSERS = {
    ".md": _parse_text,
    ".txt": _parse_text,
    ".docx": _parse_docx,
    ".pptx": _parse_pptx,
    ".xlsx": _parse_xlsx,
    ".csv": _parse_csv,
}


def parse_document_to_markdown(
    file_path: Path,
    unprocessed_pages: Optional[list] = None,
    on_page_processed: Optional[Callable[[int, int], None]] = None,
    page_stats: Optional[dict] = None,
) -> str:
    """Document -> markdown: the ingestion pipeline's first stage. Every
    format is normalized to markdown WITH real `#`/`##` headings where the
    source format has structure (Word/PowerPoint styles, PDF page breaks) --
    split_by_headings only recognizes `^#{1,6}\\s+`, so a parser returning
    flat prose would silently drop the citation-preservation guarantee
    (chunk_document's docstring) for that format.

    .md/.txt pass through unchanged. .pdf/images may invoke OCR
    (app.services.ocr) per page/whole-image when no usable embedded text
    layer exists. For .pdf, a page that needs OCR and can't get it is
    skipped rather than failing the whole document -- pass a list via
    `unprocessed_pages` to receive {"page", "reason"} entries for each one
    (see _parse_pdf's docstring for when this instead raises
    OcrUnavailableError: only when the WHOLE document produced no text).
    A whole-image upload (.png/.jpg/...) has no such partial case -- OCR
    failing there fails the only "page" there is, so OcrUnavailableError
    propagates uncaught, exactly as before.

    .doc/.ppt/.xls (legacy binary Office formats) are rejected outright --
    no parser here reads the old OLE container format, and misreading one
    as text would be worse than a clear conversion instruction.

    Every format's output passes through
    app.services.citations.normalize_arabic_text (NFKC + strip tatweel/
    harakat/bidi marks) before returning -- a PDF's table text can extract
    as Unicode Presentation Forms B glyphs instead of standard Arabic
    letters (measured on a real 80-page administrative guide: a
    presentation-form spelling of a teh-marbuta-ending word contained no
    teh marbuta codepoint at all, so citation regexes failed outright on
    it), and this is the single choke point every parser's output already
    passes through. Lossless and a no-op on non-Arabic/already-clean text,
    so applying it unconditionally is never harmful -- see that function's
    docstring for why this is safe to store/embed, unlike the lossy,
    comparison-only fold_arabic in the same module.

    `on_page_processed(page_num, total_pages)`, if given, is called once
    per PDF page right after that page's own section is appended (i.e.
    AFTER any OCR call for that page has already returned) -- this is what
    app.services.ingest_jobs.process_source_file uses to keep
    source_files.pages_done current DURING processing, so a poll while an
    OCR-heavy document is still ingesting sees real progress instead of a
    status stuck at 'processing' with no other signal for potentially
    tens of minutes. PDF-only: every other format's whole-document parse
    is fast enough (no OCR in the common case) that a live counter adds
    more overhead than value; a caller that passes this for a non-PDF
    upload gets it silently ignored, matching how `unprocessed_pages` is
    already PDF-only.

    `page_stats` is the same shape of out-param: a dict this fills with
    {"page_count", "ocr_pages"} for a PDF, so the caller doesn't have to
    re-parse the file to learn them (see _parse_pdf). PDF-only, silently
    ignored for every other format.
    """
    suffix = file_path.suffix.lower()
    if suffix in LEGACY_BINARY_FORMATS:
        raise ValueError(
            f"{file_path.name}: legacy binary {suffix} is not supported. "
            f"Re-save it as {LEGACY_BINARY_FORMATS[suffix]} and upload again."
        )
    if suffix == ".pdf":
        raw = _parse_pdf(
            file_path, unprocessed_pages=unprocessed_pages,
            on_page_processed=on_page_processed, page_stats=page_stats,
        )
    elif suffix in _IMAGE_EXTENSIONS:
        raw = _parse_image(file_path)
    else:
        parser = _PARSERS.get(suffix)
        if parser is None:
            raise ValueError(f"{file_path}: unsupported file type {suffix!r}")
        raw = parser(file_path)

    from app.services.citations import normalize_arabic_text
    return normalize_arabic_text(raw)


# Real HTML tag names strip_markdown removes -- a WHITELIST, not the
# wildcard `<[^>]+>` this replaced. That wildcard treated any `<...>` span
# as a tag, which silently deleted regulatory content that happens to use
# angle brackets as inequality signs: measured live,
# "seuil <= 5 mg/m3 selon <NF EN 166>" collapsed to "seuil" -- one stray
# `<` swallowed the rest of the sentence. Requiring a known tag name
# immediately after `<`/`</` means "temperature < 40 C" and "<NF EN 166>"
# are left untouched (neither "40" nor "NF" matches an alternative), while
# real markup (PaddleOCR-VL's HTML tables, `_parse_pptx`'s notes) is still
# removed.
_HTML_TAG_NAMES = (
    "table", "thead", "tbody", "tfoot", "tr", "td", "th",
    "p", "div", "span", "br", "hr",
    "h1", "h2", "h3", "h4", "h5", "h6",
    "ul", "ol", "li",
    "b", "i", "em", "strong", "u",
    "a", "img",
)
_HTML_TAG_RE = re.compile(
    r"</?(?:" + "|".join(_HTML_TAG_NAMES) + r")\b[^<>]*/?>", re.IGNORECASE
)


def strip_markdown(text: str) -> str:
    """Remove markdown formatting artifacts so embedding model processes clean content.

    NOTE: kept byte-for-byte in sync with
    scripts/verify_ocr_arabic.py::_strip_markdown_standalone (that copy
    exists for OCR-only venvs that can't import this module's heavy
    dependencies) -- change both together.
    """
    # Turn HTML structural boundaries into whitespace BEFORE stripping tags.
    # The generic strip below deletes tags with no replacement, which
    # FUSES the text on either side: measured on a real scanned page,
    # PaddleOCR-VL returns tables as HTML, and
    # "<td>الترخيص</td><td>450 متر</td>" collapsed to "الترخيص450 متر" --
    # two unrelated cells welded into one nonsense token, embedded and
    # retrieved that way. Markdown tables never had this problem (the
    # `\|\s*` rule further down already separates their cells); only the
    # HTML path did, which is exactly the path OCR'd tables arrive on.
    text = re.sub(r"</\s*(?:td|th)\s*>", " ", text, flags=re.IGNORECASE)
    text = re.sub(
        r"</\s*(?:tr|p|div|li|h[1-6])\s*>|<\s*br\s*/?\s*>", "\n", text, flags=re.IGNORECASE
    )
    # Remove real HTML tags only (see _HTML_TAG_RE) -- NOT every `<...>`
    # span, which used to eat inequality/threshold text.
    text = _HTML_TAG_RE.sub("", text)
    # Remove markdown headers (##, ###, etc.)
    text = re.sub(r"^#{1,6}\s+", "", text, flags=re.MULTILINE)
    # Bold/italic -- guarded on both sides so this only fires on genuine
    # markdown emphasis, not on document codes or arithmetic that happen
    # to contain '*'/'_'. `(?<!\w)...(?!\w)` requires the delimiter run
    # not be glued to a word character (so 'D_Ain_el_Abd_1970' and
    # 'ISO_45001' are never touched -- every underscore in them sits
    # between two letters/digits); `(?!\s)...(?<!\s)` requires no
    # whitespace immediately inside the delimiters (so 'L * l * h' and
    # '5 * 3 = 15', where every '*' is surrounded by spaces, are never
    # touched). '**gras**' and '_italic_' still collapse to 'gras'/'italic'.
    text = re.sub(r"(?<!\w)\*{1,3}(?!\s)([^*]+)(?<!\s)\*{1,3}(?!\w)", r"\1", text)
    text = re.sub(r"(?<!\w)_{1,3}(?!\s)([^_]+)(?<!\s)_{1,3}(?!\w)", r"\1", text)
    # Remove inline code backticks
    text = re.sub(r"`([^`]+)`", r"\1", text)
    # Blockquote marker -- only when followed by a letter, not a digit.
    # `_parse_pptx` deliberately emits "> <note text>" for speaker notes
    # (always prose, starts with a letter); a line starting "> 40 C" is a
    # threshold, not a blockquote, and the old unconditional strip
    # inverted its meaning by deleting the '>'.
    text = re.sub(r"^>\s+(?!\d)", "", text, flags=re.MULTILINE)
    # Remove horizontal rules
    text = re.sub(r"^[-*_]{3,}\s*$", "", text, flags=re.MULTILINE)
    # Markdown alignment rows only -- require 3+ dashes per cell (the
    # normal markdown convention: "---", ":---:"). A single-dash cell
    # ("| - | - |") is this corpus's "neant/sans objet" convention for a
    # real, empty-valued data row, not a table separator; the old 1+-dash
    # pattern deleted that row entirely instead of collapsing it to text
    # like every other row.
    text = re.sub(r"^\|(?:\s*:?-{3,}:?\s*\|)+\s*$", "", text, flags=re.MULTILINE)
    # Collapse UNESCAPED pipes into plain text. `_render_row` (below)
    # escapes a literal pipe inside a cell as `\|` specifically so it
    # survives this step instead of being mistaken for a cell boundary;
    # the negative lookbehind honors that escaping, and the following
    # line then turns `\|` back into a literal `|`.
    text = re.sub(r"(?<!\\)\|\s*", " ", text)
    text = text.replace("\\|", "|")
    # Links -- keep BOTH the display text and the target. In this corpus
    # the parenthetical is usually a legal/article reference, not a URL
    # ("[modifie](Loi 65-99)"), and the old display-text-only rule
    # silently deleted it.
    text = re.sub(r"\[([^\]]+)\]\(([^)]+)\)", r"\1 (\2)", text)
    # Collapse multiple blank lines
    text = re.sub(r"\n{3,}", "\n\n", text)
    # Decode HTML entities (&lt; &amp; &nbsp; ...) LAST, after every
    # markup-handling step above, so a decoded entity can never be
    # mistaken for new markup by an earlier rule.
    text = html.unescape(text).replace("\xa0", " ")
    # Strip leading/trailing whitespace
    text = text.strip()
    return text


# One psycopg2 pool per distinct database URL, for the RAW-SQL path
# (app.services.search's vector query and insert_documents) -- the
# SQLAlchemy side has its own shared pool in app.models.db.
#
# Why this exists: get_db_connection used to be a bare psycopg2.connect()
# and every caller closed the connection when it was done, so the pgvector
# path paid a full TCP connect + TLS/startup + auth round trip PER SEARCH.
# A single chat turn issues two searches (the tier-2 domain vote in
# app.services.routing.resolve_domain, then the domain-scoped retrieval in
# app.routers.chat._retrieve_context), so that was two connection setups
# on the critical path of every message, before a single vector was
# compared. Reusing a pooled connection makes that cost once-per-process.
_POOL_MIN_CONN = 1
_POOL_MAX_CONN = 12
_pools: dict = {}
_pools_lock = threading.Lock()


class _PooledConnection:
    """psycopg2 connection wrapper whose .close() RETURNS the connection to
    its pool instead of tearing it down.

    Deliberately a wrapper rather than handing out the raw connection and
    asking callers to remember putconn(): every existing call site here and
    in app.services.search is written as `conn = get_db_connection(...)` /
    `finally: conn.close()`, and silently changing what close() means for a
    raw connection object would be far easier to get wrong. `.conn` is
    exposed because insert_documents/search_similar_chunks already probe
    for exactly that attribute (`conn.conn.cursor() if hasattr(conn, "conn")`)
    to support this shape.

    On close, an in-progress transaction is rolled back before the
    connection goes back to the pool. Without that, one failed query
    (e.g. a mid-migration UndefinedColumn) would poison the connection --
    every later borrower would get `InFailedSqlTransaction` on a perfectly
    valid statement, which is the classic way a "pooling made everything
    break" incident starts.
    """

    __slots__ = ("conn", "_pool", "_closed")

    def __init__(self, conn, pool):
        self.conn = conn
        self._pool = pool
        self._closed = False

    def cursor(self, *args, **kwargs):
        return self.conn.cursor(*args, **kwargs)

    def commit(self):
        return self.conn.commit()

    def rollback(self):
        return self.conn.rollback()

    def close(self):
        if self._closed:
            return
        self._closed = True
        broken = self.conn.closed != 0
        if not broken:
            try:
                # psycopg2 exposes the server-side transaction state; an
                # INTRANS/INERROR connection must not be handed to the next
                # borrower mid-transaction.
                if self.conn.get_transaction_status() != psycopg2.extensions.TRANSACTION_STATUS_IDLE:
                    self.conn.rollback()
            except Exception:
                broken = True
        try:
            self._pool.putconn(self.conn, close=broken)
        except Exception:
            logger.exception("could not return a pooled connection; discarding it")

    def __enter__(self):
        return self

    def __exit__(self, exc_type, exc, tb):
        if exc_type is not None:
            try:
                self.conn.rollback()
            except Exception:
                pass
        self.close()
        return False


def _get_pool(database_url: str):
    pool = _pools.get(database_url)
    if pool is None:
        with _pools_lock:
            pool = _pools.get(database_url)
            if pool is None:
                pool = psycopg2.pool.ThreadedConnectionPool(
                    _POOL_MIN_CONN, _POOL_MAX_CONN, database_url,
                    # Same fail-fast reasoning as app.models.db: a chat
                    # request must not hang for the platform's default TCP
                    # timeout when Postgres is unreachable -- the caller
                    # (app.routers.chat._retrieve_context) is written to
                    # degrade to the disk corpus on an exception, but only
                    # if it actually gets one in reasonable time.
                    connect_timeout=5,
                )
                _pools[database_url] = pool
                logger.info("psycopg2 pool created (max=%d)", _POOL_MAX_CONN)
    return pool


def get_db_connection(database_url: str):
    """Borrow a pooled PostgreSQL connection.

    Returns a _PooledConnection: call .close() when done (every existing
    call site already does, in a finally:) and the underlying connection
    goes back to the pool rather than being discarded.

    Falls back to a plain, unpooled psycopg2.connect() if the pool is
    exhausted or cannot be created -- a saturated pool must degrade to the
    old behaviour (slower) rather than fail a request outright.
    """
    try:
        pool = _get_pool(database_url)
        return _PooledConnection(pool.getconn(), pool)
    except psycopg2.pool.PoolError:
        logger.warning("psycopg2 pool exhausted; falling back to a direct connection")
        return psycopg2.connect(database_url)


def close_db_pools() -> None:
    """Close every pooled connection. Called from app shutdown so a reload
    doesn't leave Postgres backends behind."""
    with _pools_lock:
        for url, pool in _pools.items():
            try:
                pool.closeall()
            except Exception:
                logger.exception("could not close psycopg2 pool for %s", url.split("@")[-1])
        _pools.clear()


# maxsize=2, not 4: each entry is a fully materialized SentenceTransformer
# (bge-m3 is ~2.2GB), so this cache's ceiling is measured in gigabytes of
# resident memory, not in entries. Two covers the only case that legitimately
# needs more than one in a process -- a migration script comparing the old
# and new embedding models (scripts/migrate_to_bge_m3.py's rollback path).
@lru_cache(maxsize=2)
def _load_embedding_model_cached(model_name: str) -> SentenceTransformer:
    return _load_embedding_model(model_name)


def load_embedding_model(model_name: str = DEFAULT_EMBEDDING_MODEL) -> SentenceTransformer:
    """Load and cache the embedding model.

    The cache lives on the private _load_embedding_model_cached, keyed on
    an EXPLICIT model_name, rather than directly on this function. With
    @lru_cache applied here, `load_embedding_model()` and
    `load_embedding_model(DEFAULT_EMBEDDING_MODEL)` are two different cache
    keys for the same model -- functools.lru_cache keys on the call
    arguments as given, and does not know that the omitted argument has
    that exact default. Both spellings are used in this codebase (app/main
    .py's preload and app.services.retrieval call the first; embed_query
    calls the second), so a process could hold TWO complete copies of a
    2.2GB model and pay the multi-second load twice. Measured live while
    benchmarking the query-embedding cache: the second load showed up as
    an 8-second "cold" embed on a model the app had already preloaded at
    startup.

    Was previously re-loading from disk on every call (docstring claimed
    caching that never existed) -- harmless for one-off ingestion scripts,
    but every chat/quiz request calls this via build_rag_context, so serving
    was paying a multi-second reload per message.

    Asserts the loaded model's actual dimension matches
    settings.embedding_dim -- documents.embedding is a fixed-width
    pgvector column (Vector(settings.embedding_dim)), so a mismatch here
    would otherwise surface as a cryptic pgvector INSERT failure deep in a
    batch instead of a clear error at the first load. Loud at boot beats
    silent corruption mid-ingest.
    """
    return _load_embedding_model_cached(model_name)


def _load_embedding_model(model_name: str) -> SentenceTransformer:
    print(f"Loading embedding model: {model_name}")
    model = SentenceTransformer(model_name)
    actual_dim = model.get_sentence_embedding_dimension()
    print(f"Model loaded. Embedding dimension: {actual_dim}")
    expected_dim = get_settings().embedding_dim
    if actual_dim != expected_dim:
        raise RuntimeError(
            f"Embedding model {model_name!r} produces {actual_dim}-dim vectors, "
            f"but settings.embedding_dim is {expected_dim}. The documents.embedding "
            f"column is a fixed-width vector({expected_dim}) -- update settings."
            f"embedding_dim (and re-run scripts/migrate_to_bge_m3.py-style migration) "
            f"or load a matching model."
        )
    return model


def chunk_text(
    text: str,
    chunk_size: int = CHUNK_SIZE,
    chunk_overlap: int = CHUNK_OVERLAP,
) -> list[str]:
    """Split text into overlapping chunks using LangChain."""
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        length_function=len,
        separators=["\n\n", "\n", ". ", " ", ""],
    )
    return splitter.split_text(text)


def split_by_headings(text: str) -> list[tuple[str, str]]:
    """Split raw (not yet markdown-stripped) text into (heading, body)
    pairs, one per `#`-`######` section. Any text before the first heading
    (title line, metadata block) is kept as its own section with an empty
    heading, so nothing is silently dropped.
    """
    matches = list(_HEADING_RE.finditer(text))
    if not matches:
        return [("", text)]

    sections: list[tuple[str, str]] = []
    if matches[0].start() > 0:
        preamble = text[: matches[0].start()].strip()
        if preamble:
            sections.append(("", preamble))

    for i, m in enumerate(matches):
        heading = m.group(2).strip()
        body_start = m.end()
        body_end = matches[i + 1].start() if i + 1 < len(matches) else len(text)
        body = text[body_start:body_end].strip()
        sections.append((heading, body))

    return sections


def chunk_document(
    text: str,
    chunk_size: int = CHUNK_SIZE,
    chunk_overlap: int = CHUNK_OVERLAP,
) -> list[dict]:
    """Heading-aware chunking: every chunk carries its own section
    heading(s) inline, so a chunk boundary can never separate a heading
    (where an article/section reference like "Art. 283" or "المادة 4"
    often lives, never repeated in the body prose) from its body. Without
    this, retrieval can return just the body half, and extract_citations
    (app/services/citations.py) never sees the reference -- citation
    injection is silently disabled for that answer.

    Short adjacent sections (this corpus's "المادة 1" .. "المادة 7"-style
    articles are frequently one or two sentences each) are PACKED together
    up to chunk_size rather than isolated one-per-chunk: measured against
    tests/data/retrieval_eval.jsonl, isolating every section regardless of
    size cost recall (many small, heading-heavy chunks embed less
    precisely than a few chunks with more surrounding context) -- see
    docs/architecture/data-and-retrieval.md's baseline. Packing keeps that
    cross-boundary context while still guaranteeing every packed section's
    own heading text sits next to its own body in the packed chunk's
    content (not just the first one's), because each piece is rendered
    "heading\n\nbody" *before* packing.

    Only a section long enough to exceed chunk_size on its own is split
    with chunk_text, each split piece re-prefixed with the heading so a
    long section's later pieces don't lose their reference either.

    Returns a list of {"content", "heading", "section_index"} dicts.
    `heading` on a packed chunk is its first section's heading (bookkeeping
    only -- every packed section's heading is already inline in `content`,
    which is what extract_citations and embedding actually see).

    A section whose body strips down to nothing (a heading with no text
    under it, or a body that is only a horizontal rule) still emits its
    HEADING alone rather than being dropped outright -- 2026-08-23 fix:
    the heading is exactly where an article/law reference lives, and the
    old `continue`-on-empty-body path discarded it along with the
    (legitimately absent) body.
    """
    sections = split_by_headings(text)
    chunks: list[dict] = []

    buffer_parts: list[str] = []
    buffer_len = 0
    buffer_first_heading = ""
    buffer_first_index = 0

    def flush():
        nonlocal buffer_parts, buffer_len, buffer_first_heading
        if buffer_parts:
            chunks.append({
                "content": "\n\n".join(buffer_parts),
                "heading": buffer_first_heading,
                "section_index": buffer_first_index,
            })
        buffer_parts, buffer_len, buffer_first_heading = [], 0, ""

    for section_index, (heading, body) in enumerate(sections):
        clean_body = strip_markdown(body)
        if not clean_body and not heading:
            # Nothing at all in this section: split_by_headings' preamble
            # branch already excludes an empty preamble, so this mainly
            # guards a heading match with no title text after it (e.g.
            # a stray "## " line) -- there is genuinely nothing to emit.
            continue
        if not clean_body:
            # A real heading with an empty body -- e.g. two consecutive
            # article headings with no text between them, or a body that
            # strip_markdown legitimately reduces to nothing (a lone
            # horizontal rule). This used to `continue` past the section
            # entirely, which silently discarded the HEADING too --
            # exactly where an article/law reference lives in this corpus
            # (chunk_document's own docstring). Emit the heading alone so
            # the reference still reaches embedding/citation extraction,
            # even though there is no body to pack it with.
            piece = heading
        else:
            piece = f"{heading}\n\n{clean_body}" if heading else clean_body

        if len(piece) > chunk_size:
            # Long enough to need its own split -- flush whatever is
            # buffered first so it isn't glued onto this section, then
            # chunk this section alone, re-prefixing every piece.
            flush()
            for body_chunk in chunk_text(clean_body, chunk_size=chunk_size, chunk_overlap=chunk_overlap):
                content = f"{heading}\n\n{body_chunk}" if heading else body_chunk
                chunks.append({"content": content, "heading": heading, "section_index": section_index})
            continue

        if buffer_parts and buffer_len + len(piece) + 2 > chunk_size:
            flush()

        if not buffer_parts:
            buffer_first_heading = heading
            buffer_first_index = section_index
        buffer_parts.append(piece)
        buffer_len += len(piece) + 2

    flush()
    return chunks


def detect_document_language(text: str) -> str:
    """Per-file language from Arabic-script character ratio -- same
    arithmetic app.services.llm.detect_query_language uses for a query,
    applied here to a whole document.

    Fixes ingest_directory previously stamping every file in a tree with
    one language default ("fr"), which silently stored Arabic-script files
    (e.g. 1.11_ar_code_travail_salama.md) as French. Committed
    language-affinity retrieval (ADR 0002 decision 5) is meaningless until
    each document's stored language is actually correct.
    """
    arabic = sum(1 for c in text if "؀" <= c <= "ۿ")
    latin = sum(1 for c in text if c.isascii() and c.isalpha())
    return "ar" if arabic > latin else "fr"


def detect_document_domain(file_path: Path, source_dir: Path) -> Optional[str]:
    """Best-effort domain from a file's path relative to source_dir,
    assuming the raw/shared/<domain>/text/*.md convention. Returns None if
    the path doesn't fit that shape -- callers should pass domain
    explicitly in that case rather than guess.
    """
    try:
        rel_parts = file_path.resolve().relative_to(source_dir.resolve()).parts
    except ValueError:
        return None
    if len(rel_parts) >= 3 and rel_parts[-2] == "text":
        domain = rel_parts[-3]
        return DOMAIN_DIR_ALIASES.get(domain, domain)
    return None


def to_vector_literal(values) -> str:
    """Format an embedding as the `[0.1,0.2,...]` text literal pgvector
    parses.

    Was `str(embedding)` on a Python list, which calls repr() on all 1024
    floats and emits full 17-significant-digit round-trip precision plus a
    space after every comma -- about 21KB per vector. Formatting at 7
    significant digits instead is lossless for the float32 the model
    actually produced (float32 carries ~7 decimal digits; the values were
    upcast to float64 by .tolist() on the way here, adding digits that
    were never real) and roughly halves both the string-building cost and
    the bytes pushed to Postgres. That matters twice: once per query on
    the search path, and once per CHUNK on the ingest path, where a
    500-chunk document was building ~10MB of Python string.
    """
    return "[" + ",".join(f"{float(v):.7g}" for v in values) + "]"


# Bounded so a long-running server can't accumulate query vectors without
# limit: 512 entries x ~4KB of formatted literal is a few MB, and chat
# queries repeat far more often than that ceiling within a session.
_QUERY_EMBED_CACHE_SIZE = 512


@lru_cache(maxsize=_QUERY_EMBED_CACHE_SIZE)
def _embed_query_cached(query: str, model_name: str) -> str:
    return to_vector_literal(load_embedding_model(model_name).encode([query])[0])


def embed_query(query: str, *, model: Optional[SentenceTransformer] = None) -> str:
    """Embed one QUERY and return it as a pgvector literal, memoized.

    A single chat turn embeds the same text twice today: app.services.
    routing.resolve_domain runs an unfiltered search to vote on the domain,
    then app.routers.chat._retrieve_context runs the domain-scoped search
    that actually produces the context. Both go through
    search_similar_chunks with the same query string, so before this the
    turn paid two full bge-m3 forward passes where one would do. The vote
    is described in resolve_domain's docstring as "sub-ms at this corpus
    size" -- true of the SQL, but the embedding in front of it is not.

    The cache is keyed on (query, model_name) so a model swap (settings.
    embedding_model, e.g. scripts/migrate_to_bge_m3.py's rollback path)
    can never serve a vector from the wrong model -- which would be
    silently wrong rather than loudly broken, since a stale 1024-dim
    vector from another model is still a valid pgvector literal.

    An explicitly-passed `model` bypasses the cache: callers that hand in
    their own SentenceTransformer (tests, scripts/eval_retrieval.py) mean
    "use exactly this object", and the cache has no way to key on it.
    """
    if model is not None and model is not _loaded_default_model():
        return to_vector_literal(model.encode([query])[0])
    return _embed_query_cached(query, DEFAULT_EMBEDDING_MODEL)


def _loaded_default_model() -> Optional[SentenceTransformer]:
    """The default model IF it is already loaded, without triggering a
    load. Lets embed_query tell "the caller handed me the same shared
    model load_embedding_model() would have returned" (cacheable) from
    "the caller handed me a different model object" (not cacheable),
    without a multi-second load as a side effect of the check.
    """
    if _load_embedding_model_cached.cache_info().currsize == 0:
        return None
    return load_embedding_model(DEFAULT_EMBEDDING_MODEL)


def embed_chunks(
    model: SentenceTransformer,
    chunks: list[str],
) -> list[list[float]]:
    """Generate embeddings for a list of text chunks.

    Batch size comes from settings.embedding_batch_size, and a CUDA
    out-of-memory error halves it and retries rather than failing the
    whole document. This is not defensive padding: this deployment runs
    the embedding model, the Ollama tutor model AND (during an upload with
    scanned pages) the resident PaddleOCR-VL worker on one 8GB card, and
    the chunk size went from 400 to 2000 characters in the bge-m3
    migration without the batch size being revisited -- so a 64-chunk
    batch of 2000-char chunks is a much larger activation spike than the
    number was originally chosen against. Losing an 80-page ingest at the
    embed step, after tens of minutes of OCR, is the expensive failure
    here; a slower retry is not.
    """
    batch_size = get_settings().embedding_batch_size
    while True:
        try:
            embeddings = model.encode(
                chunks, show_progress_bar=False, batch_size=batch_size
            )
            return embeddings.tolist()
        except Exception as e:
            if batch_size <= 1 or not _is_oom(e):
                raise
            batch_size = max(1, batch_size // 2)
            logger.warning(
                "embedding batch hit an out-of-memory error; retrying at batch_size=%d",
                batch_size,
            )
            _release_cuda_cache()


def _is_oom(exc: Exception) -> bool:
    name = type(exc).__name__
    message = str(exc).lower()
    return (
        name == "OutOfMemoryError"
        or "out of memory" in message
        or "cuda error" in message and "memory" in message
    )


def _release_cuda_cache() -> None:
    try:
        import torch

        if torch.cuda.is_available():
            torch.cuda.empty_cache()
    except Exception:
        pass


def insert_documents(
    conn,
    tenant_id: str,
    source_name: str,
    source_type: str,
    language: str,
    chunks: list[str],
    embeddings: list[list[float]],
    metadata_list: Optional[list[dict]] = None,
    domain: Optional[str] = None,
    ingest_batch_id: Optional[str] = None,
    source_file_id: Optional[str] = None,
) -> int:
    """Insert chunked documents with embeddings into pgvector.

    source_file_id: NULL (default) means "belongs to the tenant's always-on
    global corpus" -- every raw/shared CLI ingest leaves this unset, exactly
    today's behaviour. A uuid string ties every chunk back to a
    source_files row (app/routers/ingest.py's upload pipeline), which is
    what lets app.services.search's source_ids filter and the
    enable/disable toggle work per-upload.
    """
    if not metadata_list:
        metadata_list = [{} for _ in chunks]

    rows = []
    for chunk, embedding, meta in zip(chunks, embeddings, metadata_list):
        doc_id = str(uuid.uuid4())
        rows.append((
            doc_id,
            tenant_id,
            chunk,
            source_name,
            source_type,
            domain,
            ingest_batch_id,
            language,
            # to_vector_literal, not str(embedding): repr() of a 1024-float
            # Python list is ~21KB per chunk at a precision the float32
            # model never produced. See that function.
            to_vector_literal(embedding),
            json.dumps(meta),
            source_file_id,
        ))

    insert_query = """
        INSERT INTO documents
            (id, tenant_id, content, source_name, source_type, domain, ingest_batch_id, language, embedding, metadata, source_file_id)
        VALUES %s
    """

    cursor = conn.conn.cursor() if hasattr(conn, "conn") else conn.cursor()
    try:
        execute_values(
            cursor,
            insert_query,
            rows,
            template="(%s::uuid, %s, %s, %s, %s, %s, %s, %s, %s::vector, %s::jsonb, %s::uuid)",
            page_size=INSERT_PAGE_SIZE,
        )
        conn.commit()
    except Exception:
        # Explicit, because connections are POOLED now (get_db_connection):
        # a failed INSERT leaves the connection in an aborted transaction,
        # and handing that back to the pool makes the NEXT borrower fail
        # with InFailedSqlTransaction on a perfectly valid statement.
        # _PooledConnection.close() also rolls back as a backstop; doing it
        # here keeps the failure attributable to the statement that caused
        # it.
        conn.rollback()
        raise
    finally:
        cursor.close()
    return len(rows)


def ingest_text(
    text: str,
    source_name: str,
    tenant_id: str,
    source_type: str = "text",
    language: str = "fr",
    domain: Optional[str] = None,
    ingest_batch_id: Optional[str] = None,
    metadata: Optional[dict] = None,
    database_url: Optional[str] = None,
    embedding_model: Optional[SentenceTransformer] = None,
    source_file_id: Optional[str] = None,
) -> dict:
    """Ingest a single text document: split into heading-aware sections →
    chunk each → embed → store. See chunk_document() for why chunking is
    heading-aware rather than a flat strip-then-split pass.

    source_file_id: forwarded to insert_documents -- None (default) for
    every raw/shared CLI ingest; a uuid string for an app/routers/ingest.py
    upload."""
    db_url = database_url or get_settings().database_url

    if embedding_model is None:
        embedding_model = load_embedding_model()

    doc_chunks = chunk_document(text)
    if not doc_chunks:
        return {"chunks_created": 0, "source_name": source_name, "status": "empty"}

    chunk_texts = [c["content"] for c in doc_chunks]
    embeddings = embed_chunks(embedding_model, chunk_texts)

    base_meta = metadata.copy() if metadata else {}
    base_meta.update({
        "source_name": source_name,
        "source_type": source_type,
        "language": language,
    })
    metadata_list = [
        {**base_meta, "heading": c["heading"], "section_index": c["section_index"]}
        for c in doc_chunks
    ]

    conn = get_db_connection(db_url)
    try:
        count = insert_documents(
            conn, tenant_id, source_name, source_type, language, chunk_texts, embeddings,
            metadata_list, domain=domain, ingest_batch_id=ingest_batch_id,
            source_file_id=source_file_id,
        )
        return {"chunks_created": count, "source_name": source_name, "status": "success"}
    finally:
        conn.close()


def ingest_file(
    file_path: str,
    tenant_id: str,
    source_type: str = "text",
    language: Optional[str] = None,
    domain: Optional[str] = None,
    ingest_batch_id: Optional[str] = None,
    metadata: Optional[dict] = None,
    database_url: Optional[str] = None,
    embedding_model: Optional[SentenceTransformer] = None,
    source_file_id: Optional[str] = None,
    source_name: Optional[str] = None,
    text: Optional[str] = None,
) -> dict:
    """Ingest a single file. `language` omitted -> auto-detected per file
    from its own content (detect_document_language), not a caller-wide
    default.

    `source_name` omitted -> path.name (the CLI's raw/shared/<domain>/
    text/*.md convention, where the on-disk filename already IS the
    citation-worthy name). Explicit for uploads (app/routers/ingest.py):
    an uploaded file's stored_path is a server-generated
    <uuid>.<ext> (deliberately, to avoid path traversal/collisions -- see
    that router), which would otherwise become the citation shown to
    users (e.g. "3393e3ef014f4e1faf7f014aac6705ab.md" instead of
    "politique_consignation_nord.md") -- caught live by probe_upload_e2e.py.

    `text` omitted -> parsed from `file_path` here, as before. A caller
    that already parsed the file (app.services.ingest_jobs.process_
    source_file does exactly this, to get language/domain from the
    content before calling this function) should pass that result through
    instead of leaving this default -- otherwise the file is parsed AGAIN
    from scratch, and for a PDF needing OCR that means every OCR'd page
    runs through the resident worker TWICE per upload, doubling ingestion
    time for no benefit (measured live: an 80-page document with 25 OCR
    pages took ~53 minutes with the double-parse still in place). The two
    independent parses could also each hit different transient OCR
    failures (worker timeout/restart under sustained GPU load), so the
    version actually chunked-and-stored could silently differ in which
    pages succeeded from the version source_files.unprocessed_pages
    reported to the user -- a correctness bug, not just a speed one.
    """
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"File not found: {file_path}")

    if text is None:
        text = parse_document_to_markdown(path)
    resolved_language = language or detect_document_language(text)
    return ingest_text(
        text=text,
        source_name=source_name or path.name,
        tenant_id=tenant_id,
        source_type=source_type,
        language=resolved_language,
        domain=domain,
        ingest_batch_id=ingest_batch_id,
        metadata=metadata,
        database_url=database_url,
        embedding_model=embedding_model,
        source_file_id=source_file_id,
    )


def ingest_directory(
    source_dir: str,
    tenant_id: str,
    language: Optional[str] = None,
    domain: Optional[str] = None,
    database_url: Optional[str] = None,
) -> list[dict]:
    """Ingest all .txt and .md files from a directory.

    `language` is optional and now per-file by default (Arabic-script
    files are no longer silently stored as French because a caller passed
    one tree-wide default) -- pass it explicitly only to force every file
    in the tree to one language, overriding detection.

    `domain` resolves per file, in order, and is NEVER left as None:
    1. The raw/shared/<domain>/text/ path convention (detect_document_domain).
    2. This explicit `domain` argument (the --domain CLI flag), when the
       file doesn't fit that path shape.
    3. settings.default_domain, logged as an [INGEST WARNING] so a
       mis-filed document is a loud, audited fallback rather than a silent
       guess.
    A domain=NULL row is invisible to every domain-filtered query
    (`WHERE domain = %s` never matches NULL) -- confirmed live during the
    2026-08-11 E2E audit, where two untagged files had to be backfilled by
    hand after the fact. Ingestion must never reproduce that.

    One ingest_batch_id is stamped across the whole run so a later
    corpus_version check (Stage 3) can tell "the corpus changed" from
    "this row predates domain/language tracking".
    """
    source_path = Path(source_dir)
    if not source_path.is_dir():
        raise NotADirectoryError(f"Directory not found: {source_dir}")

    embedding_model = load_embedding_model()
    batch_id = uuid.uuid4().hex
    results = []

    for file_path in sorted(source_path.glob("**/*.md")) + sorted(source_path.glob("**/*.txt")):
        rel_path = file_path.relative_to(source_path)
        file_type = "markdown" if file_path.suffix == ".md" else "text"

        file_domain = detect_document_domain(file_path, source_path)
        if file_domain is None and domain is not None:
            file_domain = domain
        if file_domain is None:
            file_domain = get_settings().default_domain
            print(
                f"[INGEST WARNING] No domain provided for {rel_path}. "
                f"Falling back to default domain: {file_domain!r}"
            )
        print(f"Ingesting [{file_type}] domain={file_domain}: {rel_path}")

        result = ingest_file(
            file_path=str(file_path),
            tenant_id=tenant_id,
            source_type=file_type,
            language=language,
            domain=file_domain,
            ingest_batch_id=batch_id,
            metadata={"source_dir": source_dir, "relative_path": str(rel_path)},
            database_url=database_url,
            embedding_model=embedding_model,
        )
        results.append(result)
        print(f"  -> {result['chunks_created']} chunks created")

    total_chunks = sum(r["chunks_created"] for r in results)
    print(f"\nDone. {len(results)} files ingested, {total_chunks} total chunks. batch={batch_id}")
    return results


if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser(description="Ingest documents into pgvector")
    parser.add_argument("--source_dir", required=True, help="Directory with .txt or .md files")
    parser.add_argument("--tenant_id", required=True, help="Tenant identifier")
    parser.add_argument(
        "--language", default=None,
        help="Force this language for every file. Omit to auto-detect per file (default).",
    )
    parser.add_argument(
        "--domain", default=None,
        help=(
            "Fallback domain for files that don't fit the "
            "raw/shared/<domain>/text/ path convention. Omit to fall back "
            "further to settings.default_domain (logged as an "
            "[INGEST WARNING] when that happens)."
        ),
    )
    parser.add_argument("--database_url", default=None, help="PostgreSQL connection URL")
    args = parser.parse_args()

    results = ingest_directory(
        source_dir=args.source_dir,
        tenant_id=args.tenant_id,
        language=args.language,
        domain=args.domain,
        database_url=args.database_url,
    )
