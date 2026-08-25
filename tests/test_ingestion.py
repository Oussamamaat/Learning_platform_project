"""
Tests for the Stage-2 ingestion fixes: heading-aware chunking, per-file
language detection, and path-based domain detection. Extended 2026-08-11
with ingest_directory's domain resolution (--domain flag, tenant-default
fallback, never-NULL) -- see the "Automatic Domain Routing" plan, section 3.
Extended 2026-08-13 with multi-format parsing (.pdf/.docx/.pptx/.xlsx/.csv)
and the OCR seam (app.services.ocr) -- see the "Multi-format ingestion +
Sources panel" plan, Phase 2.
"""
from pathlib import Path
from unittest.mock import patch

import pytest

from app.services.ingestion import (
    XLSX_MAX_ROWS,
    _render_row,
    chunk_document,
    detect_document_domain,
    detect_document_language,
    ingest_directory,
    parse_document_to_markdown,
    split_by_headings,
    strip_markdown,
)
from app.services.ocr import OcrUnavailableError


# --- split_by_headings -------------------------------------------------

def test_splits_on_headings_of_any_level():
    text = "# Title\n\nintro\n\n## Section A\n\nbody a\n\n### Sub B\n\nbody b\n"
    sections = split_by_headings(text)
    assert [h for h, _ in sections] == ["Title", "Section A", "Sub B"]


def test_preamble_before_first_heading_kept_as_empty_heading_section():
    text = "some intro text\n\n## Section A\n\nbody\n"
    sections = split_by_headings(text)
    assert sections[0] == ("", "some intro text")
    assert sections[1][0] == "Section A"


def test_no_headings_returns_single_section():
    assert split_by_headings("just plain text") == [("", "just plain text")]


# --- chunk_document: the citation-preservation guarantee ----------------

def test_heading_travels_with_every_chunk_of_its_section():
    """The core bug fix: a heading carrying an article reference must never
    be separable from its own body text, regardless of chunk boundaries."""
    text = (
        "## Obligations du travailleur (Art. 283)\n\n"
        "Chaque travailleur doit prendre soin de sa securite."
    )
    chunks = chunk_document(text)
    assert len(chunks) == 1
    assert "Art. 283" in chunks[0]["content"]
    assert "Chaque travailleur" in chunks[0]["content"]


def test_short_adjacent_sections_are_packed_together():
    """Many one-sentence sections (this corpus's 'المادة 1'..'المادة 5'
    style articles) should be packed into shared chunks up to chunk_size,
    not isolated one-per-chunk -- isolating them measurably hurt recall
    (docs/architecture/data-and-retrieval.md)."""
    text = "\n\n".join(
        f"## Article {i}\n\nShort body text {i}." for i in range(1, 6)
    )
    chunks = chunk_document(text, chunk_size=400)
    assert len(chunks) < 5, "short sections should be packed, not one chunk each"
    # Every article's own heading must still be inline in *some* chunk's content.
    for i in range(1, 6):
        assert any(f"Article {i}" in c["content"] for c in chunks)


def test_long_section_is_still_split_with_heading_reprefixed():
    long_body = " ".join(["word"] * 200)  # comfortably over chunk_size
    text = f"## Long Section (Art. 99)\n\n{long_body}"
    chunks = chunk_document(text, chunk_size=100, chunk_overlap=10)
    assert len(chunks) > 1
    for c in chunks:
        assert "Art. 99" in c["content"], "every split piece must keep the heading"


def test_heading_guarantee_holds_at_the_2026_08_13_default_chunk_size():
    """CHUNK_SIZE moved 400 -> 2000 chars alongside the bge-m3 embedding
    swap (matched to its 8192-token window, vs. the old model's actual
    128-token max_seq_length). The heading-inline citation guarantee must
    survive that change unmodified -- re-asserted explicitly at a
    realistically long, multi-section document rather than relying on the
    other tests' shorter fixtures happening to stay under the new size."""
    from app.services.ingestion import CHUNK_SIZE

    assert CHUNK_SIZE == 2000
    long_body = " ".join(["mot"] * 1000)  # ~4000 chars, comfortably over 2000
    text = f"## Obligations du travailleur (Art. 283)\n\n{long_body}"
    chunks = chunk_document(text)
    assert len(chunks) > 1
    for c in chunks:
        assert "Art. 283" in c["content"]


def test_empty_document_returns_no_chunks():
    assert chunk_document("") == []


def test_no_heading_document_still_chunks():
    chunks = chunk_document("Just a paragraph with no markdown headings at all.")
    assert len(chunks) == 1
    assert chunks[0]["heading"] == ""


# --- chunk_document: 2026-08-23 fix -- a heading with an empty body used
# to be dropped ENTIRELY, taking its article/law reference with it -------

def test_heading_with_empty_body_is_still_emitted():
    """Two back-to-back article headings with no body between the first
    and second: the first used to vanish completely (heading AND its
    reference number), indistinguishable from that article never having
    existed in the source. It must now still appear somewhere in the
    corpus, even with nothing to say about it yet."""
    text = "## المادة 12\n\n## المادة 13\n\nنص المادة 13"
    chunks = chunk_document(text)
    all_content = " ".join(c["content"] for c in chunks)
    assert "المادة 12" in all_content
    assert "المادة 13" in all_content
    assert "نص المادة 13" in all_content


def test_heading_whose_body_is_only_a_horizontal_rule_keeps_the_heading():
    """strip_markdown legitimately reduces a lone '---' body to nothing --
    that's correct (a real horizontal rule is not content), but the
    heading above it ('Article 14 (Art. 283)', the actual citation) must
    survive even though its body does not."""
    text = "## Article 14 (Art. 283)\n\n---\n\n## Article 15\n\ncontenu"
    chunks = chunk_document(text)
    all_content = " ".join(c["content"] for c in chunks)
    assert "Art. 283" in all_content


def test_heading_only_section_does_not_reintroduce_fully_empty_chunks():
    """The fix must not regress the ORIGINAL guard it modifies: a section
    with neither a real heading NOR a body -- two consecutive '## ' marker
    lines with nothing between them, each a heading match with no title
    text and an empty body -- still contributes nothing."""
    chunks = chunk_document("## \n\n## \n\nreal content")
    for c in chunks:
        assert c["content"].strip() != ""
    assert any("real content" in c["content"] for c in chunks)


# --- detect_document_language -------------------------------------------

def test_detects_arabic_script_document():
    text = "المادة 283: يجب على المشغل أن يوفر معدات الوقاية الشخصية."
    assert detect_document_language(text) == "ar"


def test_detects_french_document():
    text = "L'employeur doit fournir des equipements de protection individuelle adaptes."
    assert detect_document_language(text) == "fr"


def test_mostly_latin_with_some_arabic_terms_stays_french():
    text = "Le contrat prevoit la reference suivante: Loi N 42-25, ISO 45001."
    assert detect_document_language(text) == "fr"


# --- detect_document_domain ----------------------------------------------

def test_derives_domain_from_shared_convention():
    source_dir = Path("raw/shared")
    file_path = Path("raw/shared/industrial/text/1.1_code_du_travail.md")
    assert detect_document_domain(file_path, source_dir) == "industrial"


def test_normalizes_securite_physique_alias():
    source_dir = Path("raw/shared")
    file_path = Path("raw/shared/securite_physique/text/2.1_loi_27_06.md")
    assert detect_document_domain(file_path, source_dir) == "securite"


def test_returns_none_for_path_not_matching_convention():
    source_dir = Path("raw/tenant_placeholder")
    file_path = Path("raw/tenant_placeholder/random_file.md")
    assert detect_document_domain(file_path, source_dir) is None


# --- strip_markdown: HTML table cell boundaries ---------------------------

def test_strip_markdown_separates_html_table_cells():
    """Measured on a real scanned page: PaddleOCR-VL returns tables as HTML,
    and the generic `<[^>]+>` strip deleted the tags with no replacement --
    welding adjacent cells into one nonsense token
    ("<td>الترخيص</td><td>450 متر</td>" -> "الترخيص450 متر"), which then got
    embedded and retrieved in that state. Markdown tables never had this
    problem; only the HTML path did, which is the path OCR'd tables use."""
    out = strip_markdown("<table border=1><tr><td>الترخيص</td><td>450 متر</td></tr></table>")
    assert "الترخيص 450 متر" in out
    assert "الترخيص450" not in out


def test_strip_markdown_treats_br_and_block_ends_as_line_breaks():
    assert strip_markdown("ligne1<br>ligne2") == "ligne1\nligne2"
    assert strip_markdown("<p>un</p><p>deux</p>").split() == ["un", "deux"]


def test_strip_markdown_leaves_markdown_tables_and_headings_intact():
    """The HTML fix must not disturb the existing markdown paths."""
    assert strip_markdown("| Col A | Col B |").split() == ["Col", "A", "Col", "B"]
    assert strip_markdown("# Titre\n\n**gras** ici") == "Titre\n\ngras ici"


# --- strip_markdown: 2026-08-23 hardening -- data corruption on a
# regulatory/technical corpus, found by auditing every rule against real
# content rather than just the HTML-table case above ----------------------

def test_strip_markdown_preserves_inequality_thresholds():
    """The old `<[^>]+>` wildcard treated ANY '<...>' span as a tag,
    deleting regulatory threshold clauses: measured live,
    'seuil <= 5 mg/m3 selon <NF EN 166>' collapsed to just 'seuil'. Requiring
    a whitelisted tag name immediately after '<' (see _HTML_TAG_RE) leaves
    these untouched -- 'NF' and a bare digit are never tag names."""
    assert strip_markdown("temperature < 40 C et pression > 2 bars") == \
        "temperature < 40 C et pression > 2 bars"
    assert strip_markdown("seuil <= 5 mg/m3 selon <NF EN 166>") == \
        "seuil <= 5 mg/m3 selon <NF EN 166>"


def test_strip_markdown_preserves_document_codes_with_underscores():
    """Measured in the live corpus: 'D_Ain_el_Abd_1970' (a geodetic datum
    name) became 'DAinelAbd1970' -- unfindable by vector or lexical search.
    The old emphasis regex had no notion of word boundaries; every
    underscore in an identifier sits between two word characters, which the
    new (?<!\\w)...(?!\\w) guards now require NOT be the case for a real
    emphasis match."""
    assert strip_markdown("D_Ain_el_Abd_1970") == "D_Ain_el_Abd_1970"
    assert strip_markdown("norme ISO_45001 et EN_166 requis") == \
        "norme ISO_45001 et EN_166 requis"


def test_strip_markdown_preserves_spaced_asterisk_arithmetic():
    """'surface = L * l * h' lost its multiplication signs entirely. Real
    markdown emphasis never has whitespace immediately inside its
    delimiters ('** gras **' is not bold); the new (?!\\s)...(?<!\\s) guards
    encode exactly that, so a '*' with a space on both sides is left alone."""
    assert strip_markdown("surface = L * l * h") == "surface = L * l * h"
    assert strip_markdown("5 * 3 = 15 et 2 * 4 = 8") == "5 * 3 = 15 et 2 * 4 = 8"


def test_strip_markdown_still_collapses_real_emphasis():
    """The hardening must not regress genuine markdown emphasis -- only
    guard against identifiers/arithmetic that happen to reuse '*'/'_'."""
    assert strip_markdown("**gras** ici") == "gras ici"
    assert strip_markdown("_italic_ mot") == "italic mot"


def test_strip_markdown_preserves_leading_threshold_angle():
    """The blockquote rule '^>\\s?' unconditionally deleted a leading '>',
    which INVERTS the meaning of a threshold ('> 40 C' -> '40 C'). Gating on
    'not immediately followed by a digit' distinguishes that from
    _parse_pptx's own blockquote notes, which are always prose."""
    assert strip_markdown("> 40 C au poste") == "> 40 C au poste"
    assert strip_markdown("> Remarque importante du formateur") == \
        "Remarque importante du formateur"


def test_strip_markdown_keeps_a_real_data_row_of_dashes():
    """'| - | - |' is this corpus's french-regulatory 'neant/sans objet'
    convention for an empty-valued DATA row, not a markdown alignment
    separator -- the old 1+-dash pattern deleted it outright. Requiring 3+
    dashes per cell (the normal markdown convention) fixes this while a
    real alignment row ('|---|---|') still vanishes as before."""
    out = strip_markdown("| - | - |")
    assert out.strip() != ""
    assert "-" in out

    out = strip_markdown("nom\n|---|---|\nval")
    assert "---" not in out
    assert "nom" in out and "val" in out


def test_strip_markdown_keeps_pipe_escaping_round_trip():
    """`_render_row` escapes a literal pipe inside a cell as `\\|` so it
    survives the generic pipe-collapse rule. The old unconditional
    `\\|\\s*` -> ' ' regex broke that contract, leaving a stray backslash
    in the stored text instead of the literal '|' the escaping was meant
    to preserve."""
    row = _render_row(["a|b", "c"])
    out = strip_markdown(row)
    assert "a|b" in out
    assert "\\" not in out


def test_strip_markdown_keeps_legal_reference_in_link():
    """The old link rule kept only the display text and silently deleted
    the parenthetical -- in this corpus that parenthetical is usually a
    legal/article reference, not a URL: 'Article 5 [modifie](Loi 65-99)'
    became 'Article 5 modifie', losing the citation entirely."""
    assert strip_markdown("Article 5 [modifie](Loi 65-99) applicable") == \
        "Article 5 modifie (Loi 65-99) applicable"


def test_strip_markdown_decodes_html_entities():
    """PaddleOCR-VL's HTML tables can carry literal entities (&lt; &amp;
    &gt; &nbsp;) that never had a decoding step -- they survived verbatim
    into embedded/stored text. Decoded LAST, after every markup-handling
    rule, so a decoded '<'/'>' can never be mistaken for new markup by an
    earlier step in the same call."""
    assert strip_markdown("5 &lt; x &amp; y &gt; 2") == "5 < x & y > 2"


# --- parse_document_to_markdown: passthrough + dispatch -------------------

def test_md_passes_through_unchanged(tmp_path):
    f = tmp_path / "doc.md"
    f.write_text("# Title\n\nBody text.", encoding="utf-8")
    assert parse_document_to_markdown(f) == "# Title\n\nBody text."


def test_txt_passes_through_unchanged(tmp_path):
    f = tmp_path / "doc.txt"
    f.write_text("plain text content", encoding="utf-8")
    assert parse_document_to_markdown(f) == "plain text content"


def test_unsupported_extension_raises_value_error(tmp_path):
    f = tmp_path / "doc.rtf"
    f.write_bytes(b"not a supported format")
    with pytest.raises(ValueError, match="unsupported"):
        parse_document_to_markdown(f)


@pytest.mark.parametrize("suffix,modern", [(".doc", ".docx"), (".ppt", ".pptx"), (".xls", ".xlsx")])
def test_legacy_binary_office_formats_rejected_with_actionable_message(tmp_path, suffix, modern):
    """Must fail loudly and specifically -- no parser here reads the old
    OLE container format, and silently misreading one as text would be
    worse than a clear conversion instruction."""
    f = tmp_path / f"doc{suffix}"
    f.write_bytes(b"fake legacy binary content")
    with pytest.raises(ValueError) as exc_info:
        parse_document_to_markdown(f)
    assert "not supported" in str(exc_info.value)
    assert modern in str(exc_info.value)


# --- parse_document_to_markdown: .pdf, text-layer vs. OCR fallback --------

class _FakeDict(dict):
    """Stands in for pypdf's DictionaryObject/IndirectObject: real pypdf
    objects need a .get_object() to resolve to the actual dict; a plain
    dict doesn't have one."""

    def get_object(self):
        return self


class _FakeBox:
    def __init__(self, width, height):
        self.width = width
        self.height = height


def test_pdf_text_layer_parses_without_invoking_ocr(tmp_path, monkeypatch):
    """A page with a real text layer, a font resource, and no full-page
    image must classify NATIVE and use the embedded text directly.
    Implicitly proves OCR was never invoked: the default ocr_engine is
    "none", so if this were misclassified as needing OCR, it would raise
    OcrUnavailableError instead of returning cleanly."""
    class _FakePage:
        def __init__(self):
            self.mediabox = _FakeBox(612, 792)  # US Letter, in points

        def extract_text(self):
            return "Chaque travailleur doit prendre soin de sa securite. " * 3

        def get(self, key):
            if key == "/Resources":
                return _FakeDict({"/Font": _FakeDict({"/F1": _FakeDict({})})})
            return None

    class _FakeReader:
        def __init__(self, path):
            self.pages = [_FakePage()]

    monkeypatch.setattr("pypdf.PdfReader", _FakeReader)

    f = tmp_path / "doc.pdf"
    f.write_bytes(b"%PDF-1.4 fake binary content")
    md = parse_document_to_markdown(f)
    assert "## Page 1" in md
    assert "Chaque travailleur" in md


def test_pdf_scanned_page_with_ocr_disabled_raises_actionable_error(tmp_path, monkeypatch):
    """A page app.services.pdf_classify classifies OCR_REQUIRED must raise
    OcrUnavailableError when ocr_engine="none" -- an actionable message
    (convert the file, or enable OCR), not a silent empty result and not an
    opaque NotImplementedError. The classification itself (which real pages
    hit OCR_REQUIRED) is covered by tests/test_pdf_classify.py; this test
    only proves the wiring from that decision to a loud, actionable
    document failure.

    ocr_engine is pinned to "none" here rather than inherited from the
    settings default: this test asserts the DISABLED-OCR behaviour, so
    reading the ambient default made it silently depend on what that
    default happened to be (it broke the moment the default became
    "paddleocr", and would have tried to run a real GPU OCR subprocess
    inside a unit test)."""
    from pypdf import PdfWriter
    from app.services.ocr import NullOcrEngine
    from app.services.pdf_classify import PageDecision, PageSignals, PageStrategy

    monkeypatch.setattr("app.services.ocr.get_ocr_engine", lambda: NullOcrEngine())

    writer = PdfWriter()
    writer.add_blank_page(width=200, height=200)
    f = tmp_path / "scanned.pdf"
    with open(f, "wb") as fh:
        writer.write(fh)

    forced_signals = PageSignals(
        char_count=0, page_area_in2=7.7, char_density=0.0, font_count=0,
        image_count=1, has_full_page_raster=True, full_raster_dpi=150,
        bad_char_ratio=0.0,
    )
    monkeypatch.setattr(
        "app.services.pdf_classify.classify_page",
        lambda page, text=None: PageDecision(PageStrategy.OCR_REQUIRED, forced_signals, "forced for test"),
    )

    with pytest.raises(OcrUnavailableError, match="OCR is not enabled"):
        parse_document_to_markdown(f)


def test_pdf_mixed_document_skips_only_the_ocr_required_pages(tmp_path, monkeypatch):
    """The benchmark finding this test pins: an 80-page real document
    measured 66 NATIVE pages and 4 OCR_REQUIRED ones. Under the old
    behaviour a single OCR_REQUIRED page raised past the whole document,
    discarding the 66 good pages over the 4 bad ones. With ocr_engine="none"
    this must now: (1) NOT raise, (2) keep every NATIVE page's text, (3)
    report exactly the OCR_REQUIRED page(s) via `unprocessed_pages`.

    ocr_engine is pinned to "none" (same reason as the test above): the
    behaviour under test is what happens when a page needs OCR and CANNOT
    get it, so inheriting the ambient default would make this silently
    depend on the deployment's engine setting -- and once that default
    became "paddleocr", it ran a real GPU OCR subprocess and skipped
    nothing, which is a different scenario entirely."""
    from pypdf import PdfWriter
    from app.services.ocr import NullOcrEngine
    from app.services.pdf_classify import PageDecision, PageSignals, PageStrategy

    monkeypatch.setattr("app.services.ocr.get_ocr_engine", lambda: NullOcrEngine())

    class _FakePage:
        def __init__(self, marker):
            self.marker = marker
            self.mediabox = _FakeBox(612, 792)

        def extract_text(self):
            return "" if self.marker == "scan" else f"Native content for {self.marker}"

        def get(self, key):
            return None

    class _FakeReader:
        def __init__(self, path):
            self.pages = [_FakePage("native"), _FakePage("scan"), _FakePage("native2")]

    monkeypatch.setattr("pypdf.PdfReader", _FakeReader)

    forced_signals = PageSignals(
        char_count=0, page_area_in2=7.7, char_density=0.0, font_count=0,
        image_count=1, has_full_page_raster=True, full_raster_dpi=150,
        bad_char_ratio=0.0,
    )

    def _fake_classify(page, text=None):
        if page.marker == "scan":
            return PageDecision(PageStrategy.OCR_REQUIRED, forced_signals, "forced for test")
        return PageDecision(PageStrategy.NATIVE, forced_signals, "native for test")

    monkeypatch.setattr("app.services.pdf_classify.classify_page", _fake_classify)

    f = tmp_path / "mixed.pdf"
    writer = PdfWriter()
    for _ in range(3):
        writer.add_blank_page(width=200, height=200)
    with open(f, "wb") as fh:
        writer.write(fh)

    unprocessed: list = []
    md = parse_document_to_markdown(f, unprocessed_pages=unprocessed)

    assert "Native content for native" in md
    assert "Native content for native2" in md
    assert len(unprocessed) == 1
    assert unprocessed[0]["page"] == 2
    assert unprocessed[0]["reason"] == "ocr_required"
    # `detail` is the OCR call's OWN failure message (subprocess error,
    # timeout, missing engine) -- distinct from `decision.reason` (why the
    # page was classified OCR_REQUIRED). Asserting only that it exists and
    # is non-empty, not its exact wording, which belongs to
    # app.services.ocr.NullOcrEngine and can change independently.
    assert unprocessed[0]["detail"]


def test_on_page_processed_fires_once_per_page_in_order(tmp_path, monkeypatch):
    """app.services.ingest_jobs.process_source_file uses this callback to
    keep source_files.pages_done current DURING processing -- pinned here
    at the point that actually calls it, independent of that wiring."""
    from pypdf import PdfWriter
    from app.services.ocr import NullOcrEngine
    from app.services.pdf_classify import PageDecision, PageSignals, PageStrategy

    monkeypatch.setattr("app.services.ocr.get_ocr_engine", lambda: NullOcrEngine())

    class _FakePage:
        def __init__(self, marker):
            self.marker = marker
            self.mediabox = _FakeBox(612, 792)

        def extract_text(self):
            return f"Native content for {self.marker}"

        def get(self, key):
            return None

    class _FakeReader:
        def __init__(self, path):
            self.pages = [_FakePage("a"), _FakePage("b"), _FakePage("c")]

    monkeypatch.setattr("pypdf.PdfReader", _FakeReader)
    forced_signals = PageSignals(
        char_count=10, page_area_in2=7.7, char_density=5.0, font_count=1,
        image_count=0, has_full_page_raster=False, full_raster_dpi=None,
        bad_char_ratio=0.0,
    )
    monkeypatch.setattr(
        "app.services.pdf_classify.classify_page",
        lambda page, text=None: PageDecision(PageStrategy.NATIVE, forced_signals, "native for test"),
    )

    f = tmp_path / "three_pages.pdf"
    writer = PdfWriter()
    for _ in range(3):
        writer.add_blank_page(width=200, height=200)
    with open(f, "wb") as fh:
        writer.write(fh)

    calls: list = []
    parse_document_to_markdown(f, on_page_processed=lambda page_num, total: calls.append((page_num, total)))

    assert calls == [(1, 3), (2, 3), (3, 3)]


def test_on_page_processed_exception_does_not_fail_the_parse(tmp_path, monkeypatch):
    """A progress callback is a side channel (e.g. a DB write) -- a bug in
    it must never take down an otherwise-successful parse."""
    from pypdf import PdfWriter
    from app.services.ocr import NullOcrEngine
    from app.services.pdf_classify import PageDecision, PageSignals, PageStrategy

    monkeypatch.setattr("app.services.ocr.get_ocr_engine", lambda: NullOcrEngine())

    class _FakePage:
        def extract_text(self):
            return "some native text"

        def get(self, key):
            return None

    class _FakeReader:
        def __init__(self, path):
            self.pages = [_FakePage()]

    monkeypatch.setattr("pypdf.PdfReader", _FakeReader)
    forced_signals = PageSignals(
        char_count=10, page_area_in2=7.7, char_density=5.0, font_count=1,
        image_count=0, has_full_page_raster=False, full_raster_dpi=None,
        bad_char_ratio=0.0,
    )
    monkeypatch.setattr(
        "app.services.pdf_classify.classify_page",
        lambda page, text=None: PageDecision(PageStrategy.NATIVE, forced_signals, "native for test"),
    )

    f = tmp_path / "one_page.pdf"
    writer = PdfWriter()
    writer.add_blank_page(width=200, height=200)
    with open(f, "wb") as fh:
        writer.write(fh)

    def _raising_callback(page_num, total):
        raise RuntimeError("simulated DB write failure")

    md = parse_document_to_markdown(f, on_page_processed=_raising_callback)
    assert "some native text" in md


def test_pdf_presentation_form_text_normalized_to_standard_arabic(tmp_path, monkeypatch):
    """Measured on a real 80-page administrative guide: a PDF's table text
    can extract as Unicode Presentation Forms B glyphs (U+FE70-FEFF)
    instead of standard Arabic letters. The presentation-form spelling of a
    word ending in teh marbuta contains no teh-marbuta codepoint at all --
    citations.py's `المادة\\s+(\\d+)` fails outright against it, before any
    ة/ه OCR-confusion question even applies. parse_document_to_markdown must
    normalize this at the single choke point every parser's output passes
    through, so it never reaches storage/embedding in the corrupted form."""
    from pypdf import PdfWriter
    from app.services.pdf_classify import PageDecision, PageSignals, PageStrategy

    # "نسخة" (copy/version) in Presentation Forms B -- byte-distinct from,
    # but the same word as, the standard-Arabic spelling used in assertions
    # below. Isolated codepoints (initial/medial/final/isolated forms), not
    # a ligature -- exactly what a PDF's per-glyph text layer emits.
    presentation_form = "ﻧﺴﺨﺔ"

    class _FakePage:
        mediabox = _FakeBox(612, 792)

        def extract_text(self):
            return presentation_form

        def get(self, key):
            return None

    class _FakeReader:
        def __init__(self, path):
            self.pages = [_FakePage()]

    monkeypatch.setattr("pypdf.PdfReader", _FakeReader)
    forced_signals = PageSignals(
        char_count=4, page_area_in2=93.5, char_density=1.0, font_count=1,
        image_count=0, has_full_page_raster=False, full_raster_dpi=None,
        bad_char_ratio=0.0,
    )
    monkeypatch.setattr(
        "app.services.pdf_classify.classify_page",
        lambda page, text=None: PageDecision(PageStrategy.NATIVE, forced_signals, "native for test"),
    )

    f = tmp_path / "table.pdf"
    writer = PdfWriter()
    writer.add_blank_page(width=200, height=200)
    with open(f, "wb") as fh:
        writer.write(fh)

    assert "ة" not in presentation_form  # the failure mode this test pins
    md = parse_document_to_markdown(f)
    assert "نسخة" in md
    assert presentation_form not in md


# --- parse_document_to_markdown: .docx/.pptx/.xlsx/.csv -------------------

def test_docx_headings_become_markdown_headings(tmp_path):
    """Word's own heading styles must map to real `#`/`##` markdown
    headings, not flat prose -- otherwise split_by_headings can't see the
    document's structure and the citation-preservation guarantee
    (chunk_document's docstring) doesn't extend to uploaded Word docs."""
    import docx

    d = docx.Document()
    d.add_heading("Obligations du travailleur (Art. 283)", level=2)
    d.add_paragraph("Chaque travailleur doit prendre soin de sa securite.")
    f = tmp_path / "doc.docx"
    d.save(str(f))

    md = parse_document_to_markdown(f)
    assert "## Obligations du travailleur (Art. 283)" in md
    sections = split_by_headings(md)
    assert any("Art. 283" in heading for heading, _ in sections)


def test_docx_table_cell_pipe_is_escaped(tmp_path):
    """`_parse_docx` builds table rows via `_render_row`, the same helper
    xlsx/csv already use, specifically so a literal '|' inside a cell is
    escaped rather than misread as a cell boundary -- before this it built
    rows inline and skipped that escaping, so a cell like 'A|B' corrupted
    the row instead of surviving as one cell's content."""
    import docx

    d = docx.Document()
    table = d.add_table(rows=1, cols=2)
    table.cell(0, 0).text = "A|B"
    table.cell(0, 1).text = "C"
    f = tmp_path / "doc.docx"
    d.save(str(f))

    md = parse_document_to_markdown(f)
    assert "A\\|B" in md  # escaped in the raw markdown, same as xlsx/csv
    assert "A|B" in strip_markdown(md)  # and round-trips back to a literal '|'


def test_pptx_slide_title_becomes_heading_and_notes_included(tmp_path):
    import pptx

    p = pptx.Presentation()
    slide = p.slides.add_slide(p.slide_layouts[1])
    slide.shapes.title.text = "Equipements de protection (Art. 99)"
    slide.placeholders[1].text = "Casque, gants, lunettes."
    slide.notes_slide.notes_text_frame.text = "Verifier la conformite CE."
    f = tmp_path / "doc.pptx"
    p.save(str(f))

    md = parse_document_to_markdown(f)
    assert "## Equipements de protection (Art. 99)" in md
    assert "Casque, gants, lunettes." in md
    assert "Verifier la conformite CE." in md


def test_pptx_table_shape_is_extracted_not_dropped(tmp_path):
    """A PowerPoint table is a SEPARATE shape type (GraphicFrame with
    has_table=True) from the text-frame shapes the rest of this parser
    handles -- it was previously invisible to _parse_pptx entirely:
    silently absent from the output, no error, no unprocessed_pages
    record, indistinguishable from a slide that never had a table."""
    import pptx
    from pptx.util import Inches

    p = pptx.Presentation()
    slide = p.slides.add_slide(p.slide_layouts[6])  # blank layout
    table_shape = slide.shapes.add_table(
        rows=2, cols=2, left=Inches(1), top=Inches(1), width=Inches(4), height=Inches(2)
    )
    table = table_shape.table
    table.cell(0, 0).text = "Norme"
    table.cell(0, 1).text = "Seuil"
    table.cell(1, 0).text = "EN 397"
    table.cell(1, 1).text = "450 metres"
    f = tmp_path / "doc.pptx"
    p.save(str(f))

    md = parse_document_to_markdown(f)
    assert "EN 397" in md
    assert "450 metres" in md


def test_xlsx_sheet_name_becomes_heading(tmp_path):
    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Inventaire EPI"
    ws.append(["Article", "Norme"])
    ws.append(["Casque", "EN 397"])
    f = tmp_path / "doc.xlsx"
    wb.save(str(f))

    md = parse_document_to_markdown(f)
    assert "## Inventaire EPI" in md
    assert "EN 397" in md


def test_xlsx_truncation_is_announced_not_silent(tmp_path, monkeypatch):
    """A sheet longer than the cap must say so explicitly -- silent
    truncation is unacceptable in a citation-grounded system."""
    monkeypatch.setattr("app.services.ingestion.XLSX_MAX_ROWS", 2)
    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active
    for i in range(5):
        ws.append([f"row{i}"])
    f = tmp_path / "big.xlsx"
    wb.save(str(f))

    md = parse_document_to_markdown(f)
    assert "omitted" in md
    assert "row0" in md and "row4" not in md


def test_csv_delimiter_sniffed_and_rendered(tmp_path):
    f = tmp_path / "doc.csv"
    f.write_text("Article;Norme\nCasque;EN 397\n", encoding="utf-8")
    md = parse_document_to_markdown(f)
    assert "EN 397" in md
    assert "| Casque | EN 397 |" in md


# --- ingest_directory: domain resolution, never NULL ----------------------
# ingest_file (which touches Postgres via ingest_text/insert_documents) is
# mocked so these test only the domain-resolution order: path convention ->
# --domain flag -> settings.default_domain (logged, never silently guessed).

def _fake_ingest_file(**kwargs):
    return {"chunks_created": 1, **kwargs}


def test_path_convention_wins_over_explicit_domain_flag(tmp_path, capsys):
    (tmp_path / "industrial" / "text").mkdir(parents=True)
    f = tmp_path / "industrial" / "text" / "doc.md"
    f.write_text("# Heading\n\nBody", encoding="utf-8")

    with patch("app.services.ingestion.ingest_file", side_effect=_fake_ingest_file) as mock_ingest:
        ingest_directory(source_dir=str(tmp_path), tenant_id="t1", domain="blockchain")

    assert mock_ingest.call_args.kwargs["domain"] == "industrial"
    assert "[INGEST WARNING]" not in capsys.readouterr().out


def test_explicit_domain_flag_used_when_path_does_not_match_convention(tmp_path, capsys):
    f = tmp_path / "doc.md"
    f.write_text("# Heading\n\nBody", encoding="utf-8")

    with patch("app.services.ingestion.ingest_file", side_effect=_fake_ingest_file) as mock_ingest:
        ingest_directory(source_dir=str(tmp_path), tenant_id="t1", domain="securite")

    assert mock_ingest.call_args.kwargs["domain"] == "securite"
    assert "[INGEST WARNING]" not in capsys.readouterr().out


def test_falls_back_to_tenant_default_and_warns_when_neither_available(tmp_path, capsys):
    f = tmp_path / "doc.md"
    f.write_text("# Heading\n\nBody", encoding="utf-8")

    with patch("app.services.ingestion.ingest_file", side_effect=_fake_ingest_file) as mock_ingest:
        ingest_directory(source_dir=str(tmp_path), tenant_id="t1", domain=None)

    from app.config import get_settings
    assert mock_ingest.call_args.kwargs["domain"] == get_settings().default_domain
    out = capsys.readouterr().out
    assert "[INGEST WARNING]" in out
    assert "doc.md" in out


def test_domain_is_never_none_on_the_ingest_file_call(tmp_path):
    """The never-NULL guarantee itself: regardless of path/flag/default,
    ingest_file must never be called with domain=None -- a NULL domain row
    is invisible to every domain-filtered query (WHERE domain = %s never
    matches NULL), which is the exact bug this closes (2026-08-11 E2E audit,
    two untagged files silently unreachable until backfilled by hand)."""
    f = tmp_path / "doc.md"
    f.write_text("# Heading\n\nBody", encoding="utf-8")

    with patch("app.services.ingestion.ingest_file", side_effect=_fake_ingest_file) as mock_ingest:
        ingest_directory(source_dir=str(tmp_path), tenant_id="t1", domain=None)

    assert mock_ingest.call_args.kwargs["domain"] is not None
